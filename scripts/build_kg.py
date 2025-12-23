import os, re, json
from collections import defaultdict, Counter

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

SOURCE_DIR = "./source"
OUTPUT_DIR = "./knowledge-graph"

os.makedirs(OUTPUT_DIR, exist_ok=True)

def extract_text(path):
    ext = os.path.splitext(path)[1].lower()

    if ext in [".txt", ".md", ".csv", ".json"]:
        try:
            return open(path, "r", encoding="utf-8", errors="ignore").read()
        except:
            return ""

    if ext == ".docx":
        try:
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        except:
            return ""

    if ext == ".pdf":
        try:
            reader = PdfReader(path)
            return "\n".join((page.extract_text() or "") for page in reader.pages)
        except:
            return ""

    if ext == ".pptx":
        try:
            prs = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except:
            return ""

    return ""

def find_entities(text):
    # Very simple heuristic: capitalized multi-word phrases
    candidates = re.findall(r"\b([A-Z][a-zA-Z0-9]+(?:\s+[A-Z][a-zA-Z0-9]+)*)\b", text)
    # Filter out common words
    stop = {"The", "This", "That", "And", "For", "With", "From", "Project", "Customer"}
    cleaned = [c.strip() for c in candidates if c not in stop and len(c) > 2]
    return cleaned

def find_topics(text, top_n=10):
    words = re.findall(r"\b[a-zA-Z]{5,}\b", text.lower())
    stop = {"project", "requirements", "system", "solution", "cloud", "customer", "service"}
    words = [w for w in words if w not in stop]
    freq = Counter(words)
    return [w for w, _ in freq.most_common(top_n)]

# Graph structures
nodes = {
    "file": {},
    "entity": {},
    "topic": {}
}
edges = []

entity_global_counts = Counter()
topic_global_counts = Counter()

file_texts = {}

# 1. Collect data per file
for root, _, files in os.walk(SOURCE_DIR):
    for f in files:
        path = os.path.join(root, f)
        text = extract_text(path)
        if not text.strip():
            continue

        file_id = f"file:{path}"
        nodes["file"][file_id] = {
            "id": file_id,
            "type": "file",
            "name": f,
            "path": path
        }
        file_texts[file_id] = text

        entities = find_entities(text)
        topics = find_topics(text)

        entity_counts = Counter(entities)
        topic_counts = Counter(topics)

        # Update globals
        entity_global_counts.update(entity_counts)
        topic_global_counts.update(topic_counts)

        # Create entity nodes + edges file->entity
        for e, count in entity_counts.items():
            ent_id = f"entity:{e}"
            if ent_id not in nodes["entity"]:
                nodes["entity"][ent_id] = {
                    "id": ent_id,
                    "type": "entity",
                    "name": e
                }
            edges.append({
                "type": "FILE_CONTAINS_ENTITY",
                "from": file_id,
                "to": ent_id,
                "weight": count
            })

        # Create topic nodes + edges file->topic
        for t, count in topic_counts.items():
            top_id = f"topic:{t}"
            if top_id not in nodes["topic"]:
                nodes["topic"][top_id] = {
                    "id": top_id,
                    "type": "topic",
                    "name": t
                }
            edges.append({
                "type": "FILE_CONTAINS_TOPIC",
                "from": file_id,
                "to": top_id,
                "weight": count
            })

# 2. Build entity-entity and entity-topic co-occurrence edges
for file_id, text in file_texts.items():
    entities = list({e for e in find_entities(text)})
    topics = list({t for t in find_topics(text)})

    # entity-entity
    for i in range(len(entities)):
        for j in range(i+1, len(entities)):
            e1 = f"entity:{entities[i]}"
            e2 = f"entity:{entities[j]}"
            edges.append({
                "type": "ENTITY_RELATED_TO_ENTITY",
                "from": e1,
                "to": e2,
                "file": file_id
            })

    # entity-topic
    for e in entities:
        for t in topics:
            e_id = f"entity:{e}"
            t_id = f"topic:{t}"
            edges.append({
                "type": "ENTITY_RELATED_TO_TOPIC",
                "from": e_id,
                "to": t_id,
                "file": file_id
            })

graph = {
    "nodes": {
        "file": list(nodes["file"].values()),
        "entity": list(nodes["entity"].values()),
        "topic": list(nodes["topic"].values())
    },
    "edges": edges,
    "stats": {
        "num_files": len(nodes["file"]),
        "num_entities": len(nodes["entity"]),
        "num_topics": len(nodes["topic"])
    }
}

with open(os.path.join(OUTPUT_DIR, "graph.json"), "w", encoding="utf-8") as f:
    json.dump(graph, f, indent=2)

# 3. Create Graphviz (.gv)
gv_lines = ["graph G {"]
for fnode in nodes["file"].values():
    gv_lines.append(f"  \"{fnode['id']}\" [label=\"{fnode['name']}\" shape=box];")
for enode in nodes["entity"].values():
    gv_lines.append(f"  \"{enode['id']}\" [label=\"{enode['name']}\" shape=ellipse];")
for tnode in nodes["topic"].values():
    gv_lines.append(f"  \"{tnode['id']}\" [label=\"{tnode['name']}\" shape=diamond];")

for e in edges:
    gv_lines.append(f"  \"{e['from']}\" -- \"{e['to']}\" [label=\"{e['type']}\"];")

gv_lines.append("}")
with open(os.path.join(OUTPUT_DIR, "graph.gv"), "w", encoding="utf-8") as f:
    f.write("\n".join(gv_lines))

# 4. Human-readable summary
summary_lines = []
summary_lines.append("# Knowledge Graph Summary\n")
summary_lines.append(f"- Files: {graph['stats']['num_files']}")
summary_lines.append(f"- Entities: {graph['stats']['num_entities']}")
summary_lines.append(f"- Topics: {graph['stats']['num_topics']}\n")

summary_lines.append("## Top Entities\n")
for ent, cnt in entity_global_counts.most_common(20):
    summary_lines.append(f"- {ent} ({cnt} mentions)")

summary_lines.append("\n## Top Topics\n")
for top, cnt in topic_global_counts.most_common(20):
    summary_lines.append(f"- {top} ({cnt} occurrences)")

with open(os.path.join(OUTPUT_DIR, "graph.md"), "w", encoding="utf-8") as f:
    f.write("\n".join(summary_lines))

print("Knowledge graph built.")
